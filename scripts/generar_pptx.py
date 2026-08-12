# -*- coding: utf-8 -*-
"""
KG ACADEMY - Generador del documento tecnico en PowerPoint
Infraestructura + modelo de datos (todas las tablas).

Uso:  python scripts/generar_pptx.py
Salida: docs/KG_Academy_Infraestructura_y_Base_de_Datos.pptx

Autor del desarrollo: Diego Alejandro Hernandez Blanco
"""
import os
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from kg_ppt import (
    Deck, rect, text, shape_text, table, LOGO,
    NAVY, NAVY_L, LIME, LIME_L, WHITE, GREY, CLOUD, LINE, W, H,
)

deck = Deck()

# =====================================================================
# 1. PORTADA
# =====================================================================
s = deck.blank()
rect(s, 0, 0, W, H, fill=NAVY, shape=MSO_SHAPE.RECTANGLE)
rect(s, 0, 0, Inches(0.22), H, fill=LIME, shape=MSO_SHAPE.RECTANGLE)
rect(s, W - Inches(5.0), 0, Inches(5.0), H, fill=NAVY_L, shape=MSO_SHAPE.RECTANGLE)
if os.path.exists(LOGO):
    s.shapes.add_picture(LOGO, W - Inches(4.35), Inches(1.55), height=Inches(4.4))

text(s, Inches(0.95), Inches(1.5), Inches(7.4), Inches(0.4),
     [("DOCUMENTO TECNICO  ·  VERSION 1.0", 12, True, LIME, 0)])
text(s, Inches(0.95), Inches(2.05), Inches(7.4), Inches(2.0),
     [("KG ACADEMY", 46, True, WHITE, 6),
      ("Infraestructura y base de datos", 27, True, LIME, 0)])
rect(s, Inches(0.95), Inches(4.05), Inches(1.5), Pt(4), fill=LIME, shape=MSO_SHAPE.RECTANGLE)
text(s, Inches(0.95), Inches(4.4), Inches(7.0), Inches(1.2),
     [("Plataforma educativa LMS con capacidades B2C y B2B para la gestion, "
       "trazabilidad y certificacion de la formacion en Seguridad y Salud en el Trabajo.",
       13, False, RGBColor(0xC5, 0xD5, 0xE5), 0)], spacing=1.25)
text(s, Inches(0.95), Inches(5.85), Inches(7.4), Inches(1.1),
     [("KG GESTION INTEGRAL S.A.S.  ·  KATERINE GUANARITA", 12, True, WHITE, 4),
      ("Realizado por Diego Alejandro Hernandez Blanco", 12, True, LIME, 4),
      ("Bogota D.C., Colombia  ·  2026", 10, False, GREY, 0)])

# =====================================================================
# 2. AGENDA
# =====================================================================
s = deck.page("Contenido", "Agenda del documento")
items = [
    ("01", "Vision general y alcance", "Que resuelve la plataforma y que incluye el MVP"),
    ("02", "Arquitectura de la solucion", "Capas, responsabilidades y flujo de una peticion"),
    ("03", "Stack tecnologico", "Tecnologias elegidas y justificacion"),
    ("04", "Infraestructura de despliegue", "Ambiente local y ambiente productivo en Hostinger"),
    ("05", "Mapa del modelo de datos", "11 dominios funcionales, 44 tablas"),
    ("06", "Diccionario de tablas", "Detalle campo a campo por dominio"),
    ("07", "Motor de progreso y certificacion", "Reglas de negocio implementadas"),
    ("08", "Seguridad, permisos y auditoria", "Control de acceso y trazabilidad"),
    ("09", "Roadmap por fases", "Que esta listo y que sigue"),
]
y = 1.72
for i, (n, t, d) in enumerate(items):
    col = 0 if i < 5 else 1
    yy = 1.72 + (i if i < 5 else i - 5) * 1.02
    x = Inches(0.62 + col * 6.3)
    rect(s, x, Inches(yy), Inches(6.0), Inches(0.86), fill=CLOUD, line=LINE)
    b = rect(s, x + Inches(0.14), Inches(yy + 0.14), Inches(0.58), Inches(0.58), fill=NAVY)
    shape_text(b, n, 13, True, LIME)
    text(s, x + Inches(0.88), Inches(yy + 0.16), Inches(5.0), Inches(0.6),
         [(t, 12.5, True, NAVY, 1), (d, 9.5, False, GREY, 0)])

# =====================================================================
# 3. VISION GENERAL
# =====================================================================
s = deck.page("01 · Vision general", "Que es KG Academy",
         "Plataforma tipo LMS/marketplace que combina formacion virtual con gestion corporativa de la capacitacion.")
bloques = [
    ("B2C", "Estudiante independiente", "Se registra, compra o accede a cursos, estudia, se evalua y descarga su certificado.", LIME),
    ("B2B", "Empresa cliente", "Carga trabajadores, asigna cursos con fecha limite y consulta cumplimiento por area, cargo y sede.", NAVY_L),
    ("KG", "Administracion", "Publica cursos y contenidos, administra usuarios, empresas, evaluaciones, certificados y reportes.", NAVY),
]
for i, (tag, tit, desc, color) in enumerate(bloques):
    x = Inches(0.62 + i * 4.15)
    rect(s, x, Inches(1.75), Inches(3.85), Inches(2.5), fill=WHITE, line=LINE)
    rect(s, x, Inches(1.75), Inches(3.85), Inches(0.62), fill=color)
    text(s, x + Inches(0.22), Inches(1.9), Inches(3.4), Inches(0.35), [(tag + "  ·  " + tit, 12.5, True, WHITE, 0)])
    text(s, x + Inches(0.22), Inches(2.6), Inches(3.4), Inches(1.4), [(desc, 11, False, RGBColor(0x33, 0x44, 0x55), 0)], spacing=1.2)

rect(s, Inches(0.62), Inches(4.5), Inches(12.1), Inches(2.25), fill=CLOUD, line=LINE)
text(s, Inches(0.95), Inches(4.72), Inches(11.4), Inches(0.3), [("ALCANCE DEL MVP IMPLEMENTADO", 10.5, True, LIME, 0)])
mvp = [
    "Autenticacion, registro y gestion de sesiones",
    "Roles, permisos y separacion por empresa",
    "Catalogo, ficha de curso y matricula",
    "Aula virtual con progreso leccion a leccion",
    "Evaluaciones con calificacion automatica",
    "Certificados con codigo unico y QR",
    "Panel empresarial con asignacion masiva",
    "Panel SuperAdmin y constructor de cursos",
    "Reportes exportables a CSV",
    "Notificaciones internas",
    "Gamificacion: puntos, insignias y rachas",
    "Auditoria de acciones relevantes",
]
for i, m in enumerate(mvp):
    x = Inches(0.95 + (i % 3) * 3.95)
    yy = Inches(5.15 + (i // 3) * 0.38)
    d = rect(s, x, yy + Inches(0.06), Inches(0.1), Inches(0.1), fill=LIME, shape=MSO_SHAPE.OVAL)
    text(s, x + Inches(0.22), yy, Inches(3.6), Inches(0.3), [(m, 10, False, RGBColor(0x33, 0x44, 0x55), 0)])

# =====================================================================
# 4. ARQUITECTURA
# =====================================================================
s = deck.page("02 · Arquitectura", "Arquitectura de la solucion",
         "Aplicacion full-stack monolitica modular: el mismo proyecto sirve la interfaz y la API, sobre una unica base de datos relacional.")
capas = [
    ("CAPA DE PRESENTACION", "React 18 · Server Components · Tailwind CSS",
     "Landing · Catalogo · Aula virtual · Panel empresarial · Panel SuperAdmin · Verificacion publica", LIME),
    ("CAPA DE APLICACION", "Next.js 14 App Router · Route Handlers (API REST)",
     "/api/auth · /api/aula · /api/empresa · /api/admin  —  validacion con Zod en cada endpoint", NAVY_L),
    ("CAPA DE DOMINIO", "Servicios de negocio en TypeScript",
     "Motor de progreso · Calificacion de evaluaciones · Emision de certificados · RBAC · Auditoria", NAVY),
    ("CAPA DE DATOS", "Prisma ORM 5",
     "44 tablas · migraciones versionadas · SQLite en local / PostgreSQL en produccion", RGBColor(0x07, 0x22, 0x3A)),
]
y = 1.78
for tit, tec, det, color in capas:
    rect(s, Inches(0.62), Inches(y), Inches(12.1), Inches(1.05), fill=color)
    text(s, Inches(0.95), Inches(y + 0.14), Inches(4.6), Inches(0.8),
         [(tit, 12, True, WHITE, 2), (tec, 9.5, False, RGBColor(0xC8, 0xD8, 0xE8), 0)])
    text(s, Inches(5.9), Inches(y + 0.3), Inches(6.5), Inches(0.6),
         [(det, 10, False, RGBColor(0xE4, 0xEC, 0xF4), 0)])
    y += 1.18

rect(s, Inches(0.62), Inches(6.5), Inches(12.1), Inches(0.5), fill=LIME_L, line=LIME)
text(s, Inches(0.95), Inches(6.62), Inches(11.5), Inches(0.3),
     [("Flujo de una peticion:  Navegador  →  Server Component / Route Handler  →  Servicio de dominio  →  "
       "Prisma  →  Base de datos  →  Respuesta renderizada en servidor", 9.5, True, NAVY, 0)])

# =====================================================================
# 5. STACK
# =====================================================================
s = deck.page("03 · Stack", "Stack tecnologico",
         "Tecnologias seleccionadas para un despliegue economico, mantenible y facil de entregar a KG.")
stack_rows = [
    ("Lenguaje", "TypeScript 5.6", "Tipado estatico de extremo a extremo: menos errores en produccion."),
    ("Framework", "Next.js 14 (App Router)", "Interfaz y API en un solo proyecto; renderizado en servidor."),
    ("Interfaz", "React 18 + Tailwind CSS 3", "Design system propio con la paleta extraida del logo de KG."),
    ("ORM", "Prisma 5", "Modelo unico, migraciones versionadas y consultas con tipos."),
    ("Base de datos", "SQLite (local) / PostgreSQL (produccion)", "El mismo esquema funciona en ambos motores."),
    ("Autenticacion", "JWT firmado (jose) + cookie HttpOnly", "Sesion de 8 horas, sin dependencias externas de pago."),
    ("Contrasenas", "bcryptjs (10 rondas)", "Nunca se almacena la contrasena en texto plano."),
    ("Validacion", "Zod", "Cada endpoint valida su entrada antes de tocar la base de datos."),
    ("Certificados / QR", "qrcode + impresion a PDF del navegador", "Sin licencias adicionales; QR embebido en base64."),
    ("Contenido de cursos", "Embebido (video, PDF, Genially, SCORM)", "El material vive en el proveedor; la plataforma lo referencia."),
]
table(s, Inches(0.62), Inches(1.75), Inches(12.1),
      ["Componente", "Tecnologia", "Por que se eligio"], stack_rows, [2.2, 3.4, 6.5], row_h=0.42, size=10)

# =====================================================================
# 6. INFRAESTRUCTURA
# =====================================================================
s = deck.page("04 · Infraestructura", "Ambientes de despliegue",
         "Ambiente local para revision y ambiente productivo sobre Hostinger, segun el requisito tecnico del esqueleto funcional.")

# --- Local
rect(s, Inches(0.62), Inches(1.72), Inches(5.85), Inches(4.9), fill=WHITE, line=LINE)
rect(s, Inches(0.62), Inches(1.72), Inches(5.85), Inches(0.55), fill=NAVY_L)
text(s, Inches(0.85), Inches(1.85), Inches(5.4), Inches(0.3), [("AMBIENTE LOCAL  ·  DESARROLLO Y REVISION", 11, True, WHITE, 0)])
local = [
    ("Equipo del desarrollador / KG", "Windows 11 · Node.js 20+"),
    ("Servidor de aplicacion", "next dev  —  http://localhost:3000"),
    ("Base de datos", "SQLite  —  prisma/kg_academy.db (archivo)"),
    ("Archivos estaticos", "/public  —  logotipo y recursos de marca"),
    ("Puesta en marcha", "npm install  →  npm run setup  →  npm run dev"),
]
y = 2.5
for t, d in local:
    rect(s, Inches(0.9), Inches(y), Inches(5.3), Inches(0.72), fill=CLOUD, line=LINE)
    text(s, Inches(1.1), Inches(y + 0.11), Inches(4.9), Inches(0.55),
         [(t, 10.5, True, NAVY, 1), (d, 9, False, GREY, 0)])
    y += 0.84

# --- Produccion
rect(s, Inches(6.87), Inches(1.72), Inches(5.85), Inches(4.9), fill=WHITE, line=LINE)
rect(s, Inches(6.87), Inches(1.72), Inches(5.85), Inches(0.55), fill=NAVY)
text(s, Inches(7.1), Inches(1.85), Inches(5.4), Inches(0.3), [("AMBIENTE PRODUCTIVO  ·  HOSTINGER", 11, True, WHITE, 0)])
prod = [
    ("Dominio y SSL", "kgacademy.co  —  certificado HTTPS de Hostinger"),
    ("Servidor de aplicacion", "VPS Node.js  —  next start tras PM2 / systemd"),
    ("Base de datos", "PostgreSQL 15 gestionado  —  respaldo diario"),
    ("Almacenamiento de contenido", "Video y PDF en proveedor externo, embebidos por URL"),
    ("Correo saliente", "SMTP  —  bienvenida, recordatorios y certificados"),
]
y = 2.5
for t, d in prod:
    rect(s, Inches(7.15), Inches(y), Inches(5.3), Inches(0.72), fill=CLOUD, line=LINE)
    text(s, Inches(7.35), Inches(y + 0.11), Inches(4.9), Inches(0.55),
         [(t, 10.5, True, NAVY, 1), (d, 9, False, GREY, 0)])
    y += 0.84

text(s, Inches(0.62), Inches(6.75), Inches(12.1), Inches(0.3),
     [("Migrar de local a produccion solo requiere cambiar DATABASE_URL y el provider de Prisma a \"postgresql\". "
       "El resto del codigo permanece identico.", 9.5, True, NAVY, 0)], align=PP_ALIGN.CENTER)

# =====================================================================
# 7. MAPA DEL MODELO DE DATOS
# =====================================================================
s = deck.page("05 · Modelo de datos", "Mapa de dominios",
         "44 tablas organizadas en 11 dominios funcionales. El diccionario detallado de cada dominio viene en las siguientes laminas.")
dominios = [
    ("1. Identidad y acceso", 5, "roles · permissions · role_permissions · users · sessions · password_reset_tokens", NAVY),
    ("2. Empresas B2B", 5, "companies · company_locations · areas · positions · company_members", NAVY_L),
    ("3. Planes comerciales", 2, "plans · company_subscriptions", NAVY_L),
    ("4. Catalogo de cursos", 5, "categories · courses · modules · lessons · lesson_resources", LIME),
    ("5. Matricula y progreso", 3, "enrollments · module_progress · lesson_progress", LIME),
    ("6. Evaluaciones", 6, "question_banks · questions · question_options · assessments · assessment_questions · assessment_attempts · attempt_answers", NAVY),
    ("7. Certificados", 2, "certificate_templates · certificates", LIME),
    ("8. Asignacion empresarial", 2, "assignment_batches · course_assignments", NAVY_L),
    ("9. Pagos", 4, "coupons · orders · order_items", NAVY_L),
    ("10. Notificaciones y gamificacion", 6, "notification_templates · notifications · badges · user_badges · points_ledger · streaks", LIME),
    ("11. Sistema", 2, "audit_logs · system_settings", NAVY),
]
for i, (nombre, n, tablas, color) in enumerate(dominios):
    col, row = i % 3, i // 3
    x = Inches(0.62 + col * 4.08)
    y = Inches(1.75 + row * 1.32)
    rect(s, x, y, Inches(3.85), Inches(1.16), fill=WHITE, line=LINE)
    rect(s, x, y, Inches(0.09), Inches(1.16), fill=color, shape=MSO_SHAPE.RECTANGLE)
    text(s, x + Inches(0.24), y + Inches(0.12), Inches(3.45), Inches(0.9),
         [(nombre, 11, True, NAVY, 2), (tablas, 7.5, False, GREY, 0)], spacing=1.15)

rect(s, Inches(8.78), Inches(5.71), Inches(3.94), Inches(1.16), fill=NAVY)
text(s, Inches(8.98), Inches(5.85), Inches(3.6), Inches(0.9),
     [("44 TABLAS  ·  11 DOMINIOS", 13, True, LIME, 3),
      ("Integridad referencial completa, borrado en cascada controlado e indices en las consultas criticas.",
       8.5, False, RGBColor(0xC8, 0xD8, 0xE8), 0)], spacing=1.1)

# =====================================================================
# 8+. DICCIONARIO DE TABLAS
# =====================================================================
DICCIONARIO = [
    ("1. Identidad y acceso", "Quien entra a la plataforma y que puede hacer", [
        ("roles", "Roles del sistema", "code, name, scope, isSystem", "1:N users · N:M permissions"),
        ("permissions", "Catalogo granular de permisos (modulo.accion)", "code, module, action", "N:M roles"),
        ("role_permissions", "Matriz de permisos por rol", "roleId, permissionId", "Union roles-permissions"),
        ("users", "Usuario unico (estudiante, trabajador, instructor o staff KG)",
         "email, passwordHash, documentNumber, status, roleId, companyId, lastLoginAt", "N:1 roles · N:1 companies"),
        ("sessions", "Sesiones activas con IP y user agent", "userId, token, expiresAt, revokedAt", "N:1 users"),
        ("password_reset_tokens", "Tokens de recuperacion de contrasena", "userId, token, expiresAt, usedAt", "N:1 users"),
    ]),
    ("2. Empresas B2B", "Estructura organizacional del cliente corporativo", [
        ("companies", "Empresa cliente", "nit, legalName, tradeName, economicSector, arl, riskLevel, status", "1:N members, areas, positions"),
        ("company_locations", "Sedes de la empresa", "companyId, name, city, address", "N:1 companies"),
        ("areas", "Areas o departamentos", "companyId, name, code", "N:1 companies"),
        ("positions", "Cargos con su nivel de riesgo SST", "companyId, name, riskLevel", "N:1 companies"),
        ("company_members", "Vinculo usuario-empresa con su ubicacion organizacional",
         "companyId, userId, areaId, positionId, locationId, employeeCode, status", "N:1 users, companies, areas"),
        ("plans", "Planes comerciales B2B", "code, name, maxUsers, pricePerMonth, features", "1:N subscriptions"),
        ("company_subscriptions", "Plan contratado por cada empresa", "companyId, planId, seats, startsAt, endsAt, status", "N:1 companies, plans"),
    ]),
    ("3. Catalogo de cursos", "Estructura pedagogica y contenedores de contenido", [
        ("categories", "Categorias y subcategorias del catalogo", "slug, name, icon, color, parentId, order", "Auto-relacion · 1:N courses"),
        ("courses", "Ficha completa del curso y sus reglas de negocio",
         "code, slug, title, level, durationHours, price, status, progressRule, minPassingScore, maxAttempts, "
         "requiresFinalExam, certificateValidityMonths", "N:1 categories · 1:N modules"),
        ("modules", "Modulos del curso con su peso en el progreso", "courseId, title, order, weight, isRequired, isPublished", "N:1 courses · 1:N lessons"),
        ("lessons", "Leccion y su contenedor de contenido",
         "moduleId, title, order, contentType, contentUrl, contentBody, durationMin, completionRule, weight",
         "N:1 modules · 1:N resources"),
        ("lesson_resources", "Anexos descargables de la leccion", "lessonId, title, type, url, sizeKb", "N:1 lessons"),
    ]),
    ("4. Matricula, progreso y trazabilidad", "El nucleo del seguimiento academico", [
        ("enrollments", "Matricula del usuario en el curso",
         "userId, courseId, origin, assignmentId, status, progress, startedAt, completedAt, finalScore, timeSpentSec",
         "Unica por usuario+curso"),
        ("module_progress", "Estado de avance de cada modulo", "enrollmentId, moduleId, status, progress, completedAt", "N:1 enrollments, modules"),
        ("lesson_progress", "Avance detallado leccion a leccion",
         "enrollmentId, lessonId, status, percent, timeSpentSec, lastPositionSec, views, completedAt",
         "Soporta \"continuar donde quedo\""),
        ("assignment_batches", "Lote de asignacion masiva para auditoria", "companyId, courseId, createdById, dueDate, totalTargets", "1:N assignments"),
        ("course_assignments", "Curso asignado por la empresa a un trabajador",
         "companyId, courseId, userId, batchId, isMandatory, dueDate, status, notifiedAt", "1:1 enrollment"),
    ]),
    ("5. Evaluaciones", "Banco de preguntas, examenes e intentos", [
        ("question_banks", "Banco de preguntas por curso o tema", "name, courseId, topic", "1:N questions"),
        ("questions", "Pregunta con su retroalimentacion", "bankId, type, statement, explanation, difficulty, points", "1:N options"),
        ("question_options", "Opciones de respuesta", "questionId, text, isCorrect, feedback, order", "N:1 questions"),
        ("assessments", "Evaluacion diagnostica, por modulo o final",
         "courseId, moduleId, type, minScore, maxAttempts, timeLimitMin, shuffleQuestions, showFeedback", "1:N attempts"),
        ("assessment_questions", "Preguntas que componen cada evaluacion", "assessmentId, questionId, order, points", "Union"),
        ("assessment_attempts", "Intento presentado por el estudiante",
         "assessmentId, enrollmentId, userId, attemptNo, score, correctCount, passed, durationSec", "1:N answers"),
        ("attempt_answers", "Respuesta dada a cada pregunta", "attemptId, questionId, optionId, isCorrect, points", "N:1 attempts"),
    ]),
    ("6. Certificados", "Emision, verificacion y revocacion", [
        ("certificate_templates", "Plantilla con la identidad de KG",
         "name, orientation, signerName, signerTitle, bodyTemplate, isDefault", "1:N certificates"),
        ("certificates", "Certificado emitido con datos congelados",
         "code, userId, courseId, enrollmentId, studentName, studentDocument, courseTitle, hours, finalScore, "
         "issuedAt, expiresAt, verifyUrl, qrDataUrl, status, revokedReason, revokedById, downloads",
         "1:1 enrollment · codigo unico"),
    ]),
    ("7. Comercial, notificaciones y gamificacion", "Modulos de soporte al negocio y a la experiencia", [
        ("coupons", "Cupones de descuento", "code, discountType, discountValue, maxUses, endsAt", "1:N orders"),
        ("orders", "Orden de compra B2C o B2B", "reference, userId, companyId, subtotal, discount, total, status, gateway, paidAt", "1:N items"),
        ("order_items", "Detalle de la orden", "orderId, courseId, planId, concept, quantity, unitPrice", "N:1 orders"),
        ("notification_templates", "Plantillas de aviso por canal", "code, name, channel, subject, body", "1:N notifications"),
        ("notifications", "Notificacion enviada al usuario", "userId, templateId, channel, title, message, linkUrl, isRead", "N:1 users"),
        ("badges", "Catalogo de insignias", "code, name, criteria, points", "N:M users"),
        ("user_badges", "Insignias obtenidas", "userId, badgeId, earnedAt, reason", "Union"),
        ("points_ledger", "Libro de puntos de gamificacion", "userId, points, reason, refType, refId", "N:1 users"),
        ("streaks", "Racha diaria de estudio", "userId, currentDays, longestDays, lastActiveAt", "1:1 users"),
    ]),
    ("8. Sistema", "Trazabilidad y parametrizacion", [
        ("audit_logs", "Registro de acciones relevantes con estado antes y despues",
         "userId, actorEmail, action, entity, entityId, summary, beforeJson, afterJson, ipAddress", "N:1 users"),
        ("system_settings", "Parametros configurables de la plataforma", "key, value, type, group, label", "Marca, certificados, seguridad"),
    ]),
]

for titulo, sub, filas in DICCIONARIO:
    s = deck.page("06 · Diccionario de datos", titulo, sub)
    n = len(filas)
    row_h = 0.62 if n <= 6 else (0.55 if n == 7 else 0.5)
    size = 9.5 if n <= 6 else 9
    table(s, Inches(0.62), Inches(1.75), Inches(12.1),
          ["Tabla", "Proposito", "Campos principales", "Relaciones"],
          filas, [1.85, 3.0, 5.4, 2.3], row_h=row_h, size=size)

# =====================================================================
# MOTOR DE PROGRESO
# =====================================================================
s = deck.page("07 · Reglas de negocio", "Motor de progreso y certificacion",
         "Implementado en src/lib/progress.ts. Cada regla es configurable por curso desde la tabla courses.")

pasos = [
    ("1", "Registro de la leccion", "Al abrir una leccion se crea o actualiza lesson_progress con tiempo, posicion y numero de vistas."),
    ("2", "Calculo del porcentaje", "Segun courses.progressRule: por lecciones obligatorias, por peso de lecciones o por peso de modulos."),
    ("3", "Estado de cada modulo", "module_progress pasa a no iniciado, en progreso o completado segun sus lecciones."),
    ("4", "Condicion de aprobacion", "Curso completado = 100% de lecciones obligatorias + evaluacion final aprobada (si es exigida)."),
    ("5", "Emision del certificado", "Se genera codigo unico, QR de verificacion y se congelan nombre, curso, horas y nota."),
    ("6", "Sincronizacion B2B", "La asignacion de la empresa cambia de estado y el trabajador recibe la notificacion."),
]
y = 1.78
for n, t, d in pasos:
    rect(s, Inches(0.62), Inches(y), Inches(8.2), Inches(0.72), fill=CLOUD, line=LINE)
    c = rect(s, Inches(0.78), Inches(y + 0.11), Inches(0.5), Inches(0.5), fill=LIME, shape=MSO_SHAPE.OVAL)
    shape_text(c, n, 13, True, NAVY)
    text(s, Inches(1.45), Inches(y + 0.11), Inches(7.2), Inches(0.55),
         [(t, 11, True, NAVY, 1), (d, 8.8, False, GREY, 0)])
    y += 0.82

rect(s, Inches(9.05), Inches(1.78), Inches(3.67), Inches(4.9), fill=NAVY)
text(s, Inches(9.3), Inches(2.0), Inches(3.2), Inches(0.4), [("PARAMETROS POR CURSO", 11, True, LIME, 8)])
params = [
    ("progressRule", "obligatorios / peso_lecciones / peso_modulos"),
    ("minPassingScore", "Nota minima aprobatoria (80 por defecto)"),
    ("maxAttempts", "Intentos permitidos (3 por defecto)"),
    ("requiresFinalExam", "Exige evaluacion final"),
    ("requiresAllLessons", "Exige todas las lecciones"),
    ("certificateEnabled", "Emision automatica del certificado"),
    ("certificateValidityMonths", "Vigencia del certificado en meses"),
]
y = 2.5
for k, v in params:
    text(s, Inches(9.3), Inches(y), Inches(3.2), Inches(0.55),
         [(k, 9.5, True, WHITE, 1), (v, 8, False, RGBColor(0xA9, 0xC4, 0xDD), 0)], spacing=1.1)
    y += 0.6

# =====================================================================
# SEGURIDAD
# =====================================================================
s = deck.page("08 · Seguridad", "Control de acceso, privacidad y auditoria",
         "Los permisos se validan en la interfaz y nuevamente en el servidor: ninguna respuesta depende solo del frontend.")
seg = [
    ("Autenticacion", "JWT firmado con HS256 en cookie HttpOnly, SameSite Lax, vigencia de 8 horas. Registro en la tabla sessions."),
    ("Contrasenas", "Hash bcrypt con 10 rondas. La contrasena nunca viaja ni se guarda en texto plano."),
    ("Control de acceso", "6 roles y 72 permisos granulares (modulo.accion). requireRole() protege cada layout del servidor."),
    ("Aislamiento entre empresas", "Un administrador de empresa solo consulta y modifica registros de su propio companyId; se valida en cada endpoint."),
    ("Proteccion de datos", "Consentimiento explicito de tratamiento de datos (Ley 1581 de 2012) registrado con fecha en users.acceptedDataAt."),
    ("Auditoria", "audit_logs conserva actor, accion, entidad, resumen y el estado antes y despues del cambio en JSON."),
    ("Trazabilidad academica", "El historico de lesson_progress y assessment_attempts nunca se borra por una edicion administrativa."),
    ("Verificacion publica", "Los certificados se validan sin iniciar sesion mediante codigo unico y QR."),
]
for i, (t, d) in enumerate(seg):
    col, row = i % 2, i // 2
    x = Inches(0.62 + col * 6.3)
    y = Inches(1.78 + row * 1.22)
    rect(s, x, y, Inches(6.0), Inches(1.06), fill=WHITE, line=LINE)
    rect(s, x, y, Inches(0.09), Inches(1.06), fill=LIME, shape=MSO_SHAPE.RECTANGLE)
    text(s, x + Inches(0.26), y + Inches(0.13), Inches(5.55), Inches(0.85),
         [(t, 11, True, NAVY, 2), (d, 8.8, False, GREY, 0)], spacing=1.15)

# =====================================================================
# ROADMAP
# =====================================================================
s = deck.page("09 · Roadmap", "Estado actual y siguientes fases",
         "Fase 1 entregada y funcionando. Las fases siguientes corresponden al backlog del punto 23 del esqueleto funcional.")
fases = [
    ("FASE 1  ·  ENTREGADA", LIME,
     ["Modelo de datos completo (44 tablas)", "Autenticacion, roles y permisos", "Catalogo y aula virtual",
      "Motor de progreso y evaluaciones", "Certificados con QR y verificacion", "Panel empresarial y SuperAdmin",
      "Reportes CSV y auditoria"]),
    ("FASE 2  ·  SIGUIENTE", NAVY_L,
     ["Pasarela de pagos (por definir con KG)", "Planes B2B autogestionados", "Envio de correo (SMTP)",
      "Editor visual de preguntas", "Exportacion de reportes en PDF", "Carga masiva desde Excel"]),
    ("FASE 3  ·  POSTERIOR", NAVY,
     ["Analitica avanzada y dashboards", "Automatizaciones y recordatorios", "Integracion con Power BI",
      "Firma electronica de certificados", "Login con Google / Microsoft"]),
    ("FASE 4  ·  FUTURO", RGBColor(0x07, 0x22, 0x3A),
     ["Aplicacion movil nativa", "Marketplace de instructores externos", "Integracion con HRIS empresariales",
      "API publica para clientes"]),
]
for i, (tit, color, items_f) in enumerate(fases):
    x = Inches(0.62 + i * 3.11)
    rect(s, x, Inches(1.78), Inches(2.88), Inches(4.85), fill=WHITE, line=LINE)
    shape_text(rect(s, x, Inches(1.78), Inches(2.88), Inches(0.52), fill=color), tit, 10, True, WHITE)
    y = 2.48
    for it in items_f:
        d = rect(s, x + Inches(0.18), Inches(y + 0.07), Inches(0.09), Inches(0.09), fill=color, shape=MSO_SHAPE.OVAL)
        text(s, x + Inches(0.38), Inches(y), Inches(2.35), Inches(0.55),
             [(it, 9, False, RGBColor(0x33, 0x44, 0x55), 0)], spacing=1.1)
        y += 0.56

# =====================================================================
# CIERRE
# =====================================================================
s = deck.blank()
rect(s, 0, 0, W, H, fill=NAVY, shape=MSO_SHAPE.RECTANGLE)
rect(s, 0, H - Inches(0.22), W, Inches(0.22), fill=LIME, shape=MSO_SHAPE.RECTANGLE)
if os.path.exists(LOGO):
    s.shapes.add_picture(LOGO, W / 2 - Inches(1.55), Inches(0.85), height=Inches(3.1))
text(s, Inches(1.5), Inches(4.25), W - Inches(3.0), Inches(1.6),
     [("KG ACADEMY", 34, True, WHITE, 6),
      ("Formacion que protege vidas, con la evidencia que su empresa necesita.", 15, False, LIME, 14),
      ("KG GESTION INTEGRAL S.A.S.  ·  KATERINE GUANARITA", 12, True, RGBColor(0xC5, 0xD5, 0xE5), 4)],
     align=PP_ALIGN.CENTER)
rect(s, W / 2 - Inches(3.1), Inches(6.15), Inches(6.2), Inches(0.62), fill=NAVY_L)
text(s, W / 2 - Inches(3.0), Inches(6.29), Inches(6.0), Inches(0.35),
     [("Disenado y desarrollado por Diego Alejandro Hernandez Blanco", 11.5, True, LIME, 0)], align=PP_ALIGN.CENTER)

out = deck.save("KG_Academy_Infraestructura_y_Base_de_Datos.pptx")
print("OK ->", out)
