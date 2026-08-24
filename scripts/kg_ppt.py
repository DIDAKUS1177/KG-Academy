# -*- coding: utf-8 -*-
"""
KG ACADEMY - Utilidades compartidas para generar los documentos en PowerPoint
con la identidad visual de KG Gestion Integral S.A.S.

Lo usan:  generar_pptx.py  ·  generar_manual_pptx.py  ·  generar_despliegue_pptx.py
Autor del desarrollo: Diego Alejandro Hernandez Blanco
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
LOGO = os.path.join(BASE, "public", "brand", "kg-logo.png")
DOCS = os.path.join(BASE, "docs")

# ------------------------------------------------------- paleta de marca
NAVY = RGBColor(0x0A, 0x2D, 0x4D)
NAVY_L = RGBColor(0x1B, 0x4A, 0x73)
NAVY_D = RGBColor(0x07, 0x22, 0x3A)
LIME = RGBColor(0x8F, 0xBF, 0x16)
LIME_L = RGBColor(0xE9, 0xF4, 0xC8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x6B, 0x7C, 0x8C)
CLOUD = RGBColor(0xF6, 0xF8, 0xFB)
LINE = RGBColor(0xD6, 0xE4, 0xF0)
INK = RGBColor(0x33, 0x44, 0x55)
AMBER = RGBColor(0xB4, 0x7A, 0x0C)
RED = RGBColor(0xC0, 0x39, 0x2B)

FONT = "Segoe UI"
W, H = Inches(13.333), Inches(7.5)


class Deck:
    """Envoltorio con numeración automática de laminas y cabecera de marca."""

    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width, self.prs.slide_height = W, H
        self.n = 0

    # ----------------------------------------------------------- básicos
    def blank(self):
        return self.prs.slide_layouts[6] and self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def save(self, filename):
        os.makedirs(DOCS, exist_ok=True)
        out = os.path.join(DOCS, filename)
        self.prs.save(out)
        return out

    # ----------------------------------------------------------- laminas
    def cover(self, kicker, title, subtitle, blurb):
        s = self.blank()
        rect(s, 0, 0, W, H, fill=NAVY, shape=MSO_SHAPE.RECTANGLE)
        rect(s, 0, 0, Inches(0.22), H, fill=LIME, shape=MSO_SHAPE.RECTANGLE)
        rect(s, W - Inches(5.0), 0, Inches(5.0), H, fill=NAVY_L, shape=MSO_SHAPE.RECTANGLE)
        if os.path.exists(LOGO):
            s.shapes.add_picture(LOGO, W - Inches(4.35), Inches(1.55), height=Inches(4.4))

        text(s, Inches(0.95), Inches(1.45), Inches(7.4), Inches(0.4),
             [(kicker, 12, True, LIME, 0)])
        text(s, Inches(0.95), Inches(2.0), Inches(7.4), Inches(2.0),
             [(title, 44, True, WHITE, 6), (subtitle, 25, True, LIME, 0)])
        rect(s, Inches(0.95), Inches(4.05), Inches(1.5), Pt(4), fill=LIME, shape=MSO_SHAPE.RECTANGLE)
        text(s, Inches(0.95), Inches(4.4), Inches(7.0), Inches(1.3),
             [(blurb, 13, False, RGBColor(0xC5, 0xD5, 0xE5), 0)], spacing=1.25)
        text(s, Inches(0.95), Inches(5.9), Inches(7.4), Inches(1.1),
             [("KG GESTIÓN INTEGRAL S.A.S.  ·  KATERINE GUAÑARITA", 12, True, WHITE, 4),
              ("Realizado por Diego Alejandro Hernández Blanco", 12, True, LIME, 4),
              ("Bogotá D.C., Colombia  ·  2026", 10, False, GREY, 0)])
        return s

    def page(self, eyebrow, title, sub=None):
        s = self.blank()
        rect(s, 0, 0, W, Inches(1.28), fill=NAVY, shape=MSO_SHAPE.RECTANGLE)
        rect(s, 0, Inches(1.28), W, Pt(4), fill=LIME, shape=MSO_SHAPE.RECTANGLE)
        if os.path.exists(LOGO):
            s.shapes.add_picture(LOGO, W - Inches(1.32), Inches(0.16), height=Inches(0.96))
        runs = [(eyebrow.upper(), 10.5, True, LIME, 3), (title, 25, True, WHITE, 0)]
        if sub:
            runs.append((sub, 11, False, RGBColor(0xB8, 0xC8, 0xD8), 0))
        text(s, Inches(0.62), Inches(0.2), W - Inches(2.2), Inches(1.02), runs,
             anchor=MSO_ANCHOR.MIDDLE, spacing=1.05)

        self.n += 1
        text(s, Inches(0.62), H - Inches(0.44), Inches(9), Inches(0.3),
             [("KG Academy  ·  KG Gestión Integral S.A.S.  ·  Desarrollado por Diego Alejandro Hernández Blanco",
               8.5, False, GREY, 0)])
        text(s, W - Inches(1.1), H - Inches(0.44), Inches(0.5), Inches(0.3),
             [(str(self.n), 9, True, NAVY, 0)], align=PP_ALIGN.RIGHT)
        return s

    def closing(self, claim):
        s = self.blank()
        rect(s, 0, 0, W, H, fill=NAVY, shape=MSO_SHAPE.RECTANGLE)
        rect(s, 0, H - Inches(0.22), W, Inches(0.22), fill=LIME, shape=MSO_SHAPE.RECTANGLE)
        if os.path.exists(LOGO):
            s.shapes.add_picture(LOGO, W / 2 - Inches(1.55), Inches(0.85), height=Inches(3.1))
        text(s, Inches(1.5), Inches(4.25), W - Inches(3.0), Inches(1.6),
             [("KG ACADEMY", 34, True, WHITE, 6),
              (claim, 15, False, LIME, 14),
              ("KG GESTIÓN INTEGRAL S.A.S.  ·  KATERINE GUAÑARITA", 12, True, RGBColor(0xC5, 0xD5, 0xE5), 4)],
             align=PP_ALIGN.CENTER)
        rect(s, W / 2 - Inches(3.1), Inches(6.15), Inches(6.2), Inches(0.62), fill=NAVY_L)
        text(s, W / 2 - Inches(3.0), Inches(6.29), Inches(6.0), Inches(0.35),
             [("Diseñado y desarrollado por Diego Alejandro Hernández Blanco", 11.5, True, LIME, 0)],
             align=PP_ALIGN.CENTER)
        return s


# ------------------------------------------------------------- primitivas
def rect(s, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06):
    sh = s.shapes.add_shape(shape, x, y, w, h)
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sh.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0):
    """runs = [(texto, tamaño, negrita, color, espacio_despues_pt)]"""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, (txt, size, bold, color, after) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(after)
        p.line_spacing = spacing
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = FONT
    return tb


def shape_text(sh, txt, size=11, bold=False, color=WHITE, align=PP_ALIGN.CENTER):
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = FONT
    return sh


def table(s, x, y, w, headers, rows, widths, row_h=0.32, head_h=0.36, size=9,
          head_fill=NAVY, first_col_bold=True):
    n_rows = len(rows) + 1
    gt = s.shapes.add_table(n_rows, len(headers), x, y, w,
                            Inches(head_h + row_h * len(rows))).table
    total = sum(widths)
    for i, ww in enumerate(widths):
        gt.columns[i].width = Emu(int(w * ww / total))
    gt.rows[0].height = Inches(head_h)
    for i in range(1, n_rows):
        gt.rows[i].height = Inches(row_h)

    for c, htxt in enumerate(headers):
        cell = gt.cell(0, c)
        cell.text = ""
        cell.fill.solid()
        cell.fill.fore_color.rgb = head_fill
        cell.margin_left = cell.margin_right = Inches(0.07)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        r = cell.text_frame.paragraphs[0].add_run()
        r.text = htxt
        r.font.size = Pt(size)
        r.font.bold = True
        r.font.color.rgb = WHITE
        r.font.name = FONT

    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = gt.cell(ri, ci)
            cell.text = ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 else CLOUD
            cell.margin_left = cell.margin_right = Inches(0.07)
            cell.margin_top = cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            r = cell.text_frame.paragraphs[0].add_run()
            r.text = str(val)
            r.font.size = Pt(size - 0.5)
            r.font.bold = first_col_bold and ci == 0
            r.font.color.rgb = NAVY if (first_col_bold and ci == 0) else INK
            r.font.name = FONT
    return gt


def bullet_card(s, x, y, w, h, titulo, descripcion, color=LIME, tsize=11, dsize=9):
    """Tarjeta blanca con barra de color a la izquierda."""
    rect(s, x, y, w, h, fill=WHITE, line=LINE)
    rect(s, x, y, Inches(0.09), h, fill=color, shape=MSO_SHAPE.RECTANGLE)
    text(s, x + Inches(0.26), y + Inches(0.13), w - Inches(0.5), h - Inches(0.26),
         [(titulo, tsize, True, NAVY, 2), (descripcion, dsize, False, GREY, 0)], spacing=1.15)


def step_row(s, x, y, w, numero, titulo, descripcion, h=0.72):
    """Fila numerada con circulo lima."""
    rect(s, x, y, w, Inches(h), fill=CLOUD, line=LINE)
    c = rect(s, x + Inches(0.16), y + Inches(0.11), Inches(0.5), Inches(0.5),
             fill=LIME, shape=MSO_SHAPE.OVAL)
    shape_text(c, str(numero), 13, True, NAVY)
    text(s, x + Inches(0.83), y + Inches(0.11), w - Inches(1.05), Inches(h - 0.17),
         [(titulo, 11, True, NAVY, 1), (descripcion, 8.8, False, GREY, 0)], spacing=1.12)


def note(s, x, y, w, mensaje, tono=LIME_L, borde=LIME, color=NAVY, size=9.5, h=0.5):
    rect(s, x, y, w, Inches(h), fill=tono, line=borde)
    text(s, x + Inches(0.25), y + Inches(0.12), w - Inches(0.5), Inches(h - 0.24),
         [(mensaje, size, True, color, 0)], spacing=1.15)
